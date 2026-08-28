# -*- coding: utf-8 -*-
"""V2 信函引擎（2026 CSCO 巡讲系列版式）+ 站点配置驱动。

用法:
    python build_v2.py <station.json>            # 生成全套
    python audit_v2.py <station.json>            # 配置派生审计（见 audit_v2.py）

站点 JSON 关键字段（完整样例见 examples/v2_chengdu.json）:
    template_dir / base{chair,host,hostL,spk,disc}   捐体信（顺德 V2 终稿）
    bios / portraits  简历 JSON 与照片目录
    output  输出目录
    meeting {intro,date,weekday,span,hotel,date_line}
    titles {pos,neg,her2,case}   环节讲题
    questions {kras,her2,case}   讨论问题套（config 里用 "@kras" 等引用）
    agenda_img  官方日程海报图；null = 不追加议程页（2026-08 用户决定默认不放）
    prefix {姓名: "NN-角色"}     文件命名
    letters [ ... ]  见 README 注释于本文件底部
"""
import io, os, sys, json, copy, math

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

def A(tag): return qn('a:' + tag)

CFG = {}
BASE = {}
BIOS = {}

# ---------- 通用小工具 ----------
def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in iter_shapes(sh.shapes):
                yield s

def cap_rpr(p_el):
    for r in p_el.iter(A('r')):
        rp = r.find(A('rPr'))
        if rp is not None:
            return copy.deepcopy(rp)
    return None

def rebuild_para(p_el, segments, rpr_b=None, rpr_n=None, force_rpr=None):
    for r in p_el.findall(A('r')):
        p_el.remove(r)
    for e in p_el.findall(A('endParaRPr')):
        p_el.remove(e)
    prev = p_el.find(A('pPr'))
    for text, bold in segments:
        r = p_el.makeelement(A('r'), {})
        rp = copy.deepcopy(force_rpr) if force_rpr is not None else copy.deepcopy(rpr_b if bold else rpr_n)
        if rp is None:
            rp = p_el.makeelement(A('rPr'), {})
            rp.set('lang', 'zh-CN')
        rp.set('b', '1' if bold else '0')
        r.append(rp)
        t = p_el.makeelement(A('t'), {})
        t.text = text
        r.append(t)
        if prev is not None:
            prev.addnext(r)
        else:
            p_el.append(r)
        prev = r

def set_lnspc(p_el, val):
    pPr = p_el.find(A('pPr'))
    if pPr is None:
        pPr = p_el.makeelement(A('pPr'), {})
        p_el.insert(0, pPr)
    for old in pPr.findall(A('lnSpc')):
        pPr.remove(old)
    ln = p_el.makeelement(A('lnSpc'), {})
    sp = p_el.makeelement(A('spcPct'), {})
    sp.set('val', str(val))
    ln.append(sp)
    pPr.insert(0, ln)

def bullet_pPr(p_el, pct):
    """简历行段落属性：悬挂缩进圆点，与捐体一致"""
    pPr = p_el.makeelement(A('pPr'), {})
    pPr.set('marL', '285750')
    pPr.set('indent', '-285750')
    ln = p_el.makeelement(A('lnSpc'), {})
    sp = p_el.makeelement(A('spcPct'), {})
    sp.set('val', str(pct))
    ln.append(sp)
    pPr.append(ln)
    pPr.append(p_el.makeelement(A('buFont'), {'typeface': 'Arial'}))
    pPr.append(p_el.makeelement(A('buChar'), {'char': '\u2022'}))
    p_el.insert(0, pPr)
    return pPr

def _wrap_lines(text, chars_per_line):
    total, n = 0, len(text)
    while n > 0:
        n -= min(chars_per_line, n)
        total += 1
    return max(total, 1)

def fill_bio(tf, lines, box_w_in=4.90, box_h_in=3.35):
    """简介页简历区：目标高度到姓名框上沿；每行带圆点；超长自动降字号/行距"""
    INDENT_IN = 285750 / 914400
    eff_w = box_w_in - INDENT_IN
    tiers = [(1400, 140000), (1400, 125000), (1300, 118000), (1250, 112000),
             (1200, 108000), (1150, 105000)]
    chosen = tiers[-1]
    for sz, pct in tiers:
        cpl = math.floor(eff_w * 72 / (sz / 100.0) * 0.94)
        vis = sum(_wrap_lines(l, cpl) for l in lines)
        lh_in = (sz / 100.0) / 72.0 * (pct / 100000.0)
        if vis * lh_in <= box_h_in:
            chosen = (sz, pct)
            break
    sz, pct = chosen
    rpr0 = cap_rpr(tf.paragraphs[0]._p)
    for para in list(tf.paragraphs):
        para._p.getparent().remove(para._p)
    for line in lines:
        p_el = tf._txBody.makeelement(A('p'), {})
        bullet_pPr(p_el, pct)
        r = p_el.makeelement(A('r'), {})
        rp = copy.deepcopy(rpr0) if rpr0 is not None else p_el.makeelement(A('rPr'), {})
        rp.set('lang', 'zh-CN')
        rp.set('sz', str(sz))
        r.append(rp)
        t = p_el.makeelement(A('t'), {})
        t.text = line
        r.append(t)
        p_el.append(r)
        tf._txBody.append(p_el)

def swap_photo(slide, path):
    best = None
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.image is not None:
            a = (sh.width or 0) * (sh.height or 0)
            if best is None or a > best[0]:
                best = (a, sh)
    if best is None:
        return False
    sh = best[1]
    IN = 914400
    sh.left, sh.top = int(0.35 * IN), int(4.30 * IN)
    sh.width, sh.height = int(1.98 * IN), int(2.68 * IN)
    blip = sh._element.blipFill.find(A('blip'))
    _, rid = slide.part.get_or_add_image_part(path)
    blip.set(qn('r:embed'), rid)
    return True

def rewrite_profile(slide, expert):
    info = BIOS[expert]
    tbs = [sh for sh in slide.shapes
           if sh.has_text_frame and sh.text_frame.text.strip()
           and sh.shape_type != MSO_SHAPE_TYPE.GROUP]
    bio = max(tbs, key=lambda s: len(s.text_frame.text))
    rest = [s for s in tbs if s is not bio]
    name = next((s for s in rest if '教授' in s.text_frame.text), rest[0] if rest else None)
    IN = 914400
    bio.left, bio.top = int(2.52 * IN), int(3.42 * IN)
    bio.width, bio.height = int(4.90 * IN), int(6.15 * IN)
    fill_bio(bio.text_frame, info['lines'])
    if name is not None:
        disp = expert if len(expert) >= 3 else (expert[0] + '   ' + expert[1])
        ps = name.text_frame.paragraphs
        rebuild_para(ps[0]._p, [(disp + ' 教授', True)], force_rpr=cap_rpr(ps[0]._p))
        if len(ps) > 1:
            rebuild_para(ps[1]._p, [(info['hospital'], False)], force_rpr=cap_rpr(ps[1]._p))
    if info.get('photo'):
        swap_photo(slide, info['photo'])

def rewrite_topics(slide, questions):
    label = qtb = None
    cand = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        if '讨' in t and '问' in t and '题' in t and len(t) <= 8:
            label = sh
        else:
            cand.append(sh)
    if not cand:
        raise RuntimeError('话题页识别失败')
    qtb = max(cand, key=lambda s: len(s.text_frame.text))
    IN = 914400
    qtb.left, qtb.top = int(0.42 * IN), int(4.20 * IN)
    qtb.width, qtb.height = int(7.05 * IN), int(3.10 * IN)
    tf = qtb.text_frame
    rpr0 = cap_rpr(tf.paragraphs[0]._p)
    for para in list(tf.paragraphs):
        para._p.getparent().remove(para._p)
    for q in questions:
        p_el = tf._txBody.makeelement(A('p'), {})
        set_lnspc(p_el, 145000)
        r = p_el.makeelement(A('r'), {})
        rp = copy.deepcopy(rpr0) if rpr0 is not None else p_el.makeelement(A('rPr'), {})
        rp.set('lang', 'zh-CN')
        rp.set('sz', '1700')
        r.append(rp)
        t = p_el.makeelement(A('t'), {})
        t.text = q
        r.append(t)
        p_el.append(r)
        tf._txBody.append(p_el)

# ---------- 全局替换 / 称呼 / 开场白 ----------
def global_replace(prs):
    meet = CFG['meeting']
    for sl in prs.slides:
        for sh in iter_shapes(sl.shapes):
            if not getattr(sh, 'has_text_frame', False):
                continue
            for p in list(sh.text_frame.paragraphs):
                txt = p.text
                if CFG['_old_hotel'] and CFG['_old_hotel'] in txt:
                    rebuild_para(p._p, [(meet['hotel'], True)], force_rpr=cap_rpr(p._p))
                elif CFG['_old_date_mark'] and CFG['_old_date_mark'] in txt.replace(' ', ''):
                    rebuild_para(p._p, [(meet['date_line'], True)], force_rpr=cap_rpr(p._p))

def fix_salutation(prs, name):
    PAD = '   '   # 姓名前后各三个空格，随姓名一起带下划线
    for sh in prs.slides[0].shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            if p.text.strip().startswith('尊敬的'):
                runs = p.runs
                if len(runs) >= 4:
                    runs[0].text = '尊敬的'
                    runs[1].text = PAD
                    runs[2].text = name + PAD
                    runs[3].text = '教授：'
                    for extra in runs[4:]:
                        extra._r.getparent().remove(extra._r)
                    return True
                rp_b = rp_u = rp_n = None
                for r in runs:
                    rp = r._r.find(A('rPr'))
                    u = rp is not None and rp.get('u')
                    b = rp is not None and rp.get('b') == '1'
                    if u and rp_u is None: rp_u = rp
                    elif b and rp_b is None: rp_b = rp
                    elif rp_n is None: rp_n = rp
                rebuild_para(p._p, [('尊敬的 ', False), (name, True), ('  教授：', False)],
                             rp_b or rp_u, rp_n or rp_b)
                return True
    return False

def rewrite_intro(prs, head_role):
    """开场白单段流式重建；head_role=None 表示主席（会议主席）"""
    meet = CFG['meeting']
    THANKS = '　　衷心感谢您在繁忙的日程中拨冗出席将于'
    DATE_ALL = '%s（%s）上午%s' % (meet['date'], meet['weekday'], meet['span'])
    for sh in prs.slides[0].shapes:
        if not sh.has_text_frame:
            continue
        ps = sh.text_frame.paragraphs
        target = None
        for p in ps:
            if '衷心感谢您在繁忙的日程中拨冗出席将于' in p.text:
                target = p
                break
        if target is None:
            continue
        rb = rn = None
        for r in target.runs:
            rp = r._r.find(A('rPr'))
            b = rp is not None and rp.get('b') == '1'
            if b and rb is None: rb = rp
            if not b and rn is None: rn = rp
        role = '会议主席' if head_role is None else '会议%s嘉宾' % head_role
        segs = [(THANKS, False), (DATE_ALL, False), ('举办的', False),
                ('“%s”' % meet['intro'], True), ('并担任', False), (role, True),
                ('。以下向您简要介绍会议中您参与的内容，供您参考。', False)]
        rebuild_para(target._p, segs, rb, rn)
        for pp in ps[ps.index(target)+1:]:
            if pp.text.strip():
                pp._p.getparent().remove(pp._p)
        return True
    return False

# ---------- 任务表 ----------
def grab_rprs(task_tr):
    tc2 = task_tr.findall(A('tc'))[2]
    rb = rn = None
    for rp in tc2.iter(A('rPr')):
        if rp.get('b') == '1' and rb is None: rb = rp
        if rp.get('b') != '1' and rn is None: rn = rp
    return copy.deepcopy(rb), copy.deepcopy(rn)

def fill_table(tbl, tasks, renumber=True, header=True, heights=None):
    ns_tr = tbl.findall(A('tr'))
    body_trs = list(ns_tr[1:] if header else ns_tr)
    n = len(tasks)
    rb, rn = grab_rprs(body_trs[0])
    while len(body_trs) < n:
        cl = copy.deepcopy(body_trs[-1])
        body_trs[-1].addnext(cl)
        body_trs.append(cl)
    while len(body_trs) > n:
        tbl.remove(body_trs.pop())
    if heights:
        for tr, hv in zip(body_trs, heights):
            tr.set('h', str(hv))
    for i, (tr, task) in enumerate(zip(body_trs, tasks)):
        if renumber:
            fill_cell(tr, 0, [[(str(i + 1), False)]], rb, rn)
        fill_cell(tr, 1, [[(task['time'], False)]], rb, rn)
        fill_cell(tr, 2, task['paras'], rb, rn)

def fill_cell(tr, ci, groups, rb, rn):
    txBody = tr.findall(A('tc'))[ci].find(A('txBody'))
    paras = txBody.findall(A('p'))
    while len(paras) < len(groups):
        txBody.append(copy.deepcopy(paras[-1]))
        paras = txBody.findall(A('p'))
    for extra in paras[len(groups):]:
        txBody.remove(extra)
    for p_el, segs in zip(txBody.findall(A('p')), groups):
        rebuild_para(p_el, segs, rb, rn)

def first_tables(prs):
    return [sh.table._tbl for sh in prs.slides[0].shapes if getattr(sh, 'has_table', False)]

# ---------- 行文案组装 ----------
def _qt(titles):
    if isinstance(titles, str):
        titles = [titles]
    return '和'.join('“%s”' % t for t in titles)

def invite_row(inviter, guest, titles, role):
    segs = []
    if inviter is None:
        segs.append(('您 ', True))
        segs.append(('介绍并邀请 ', False))
        segs.append((guest + ' 教授', True))
    else:
        head = inviter if inviter == '大会主席' else inviter + ' 教授'
        segs.append((head + ' ', True))
        segs.append(('介绍并邀请 ', False))
        segs.append(('您', True))
    segs.append((' 担任%s环节的' % _qt(titles), False))
    segs.append((role, True))
    return segs

def disc_list_row(names):
    return [('您 ', True), ('介绍并邀请 ', False),
            ('、'.join(names) + ' 教授', True),
            (' 担任“讨论环节”的', False), ('学术讨论嘉宾', True)]

def handoff_row(nxt, mandates):
    return [('您 ', True), ('介绍并邀请 ', False), (nxt + ' 教授', True),
            (' 担任%s环节的' % _qt(mandates), False), ('主持嘉宾', True)]

def team_summary_row():
    return [('您 ', True), ('邀请 ', False), ('大会主席团', True), (' 进行会议总结', False)]

def drawn_invite_row(meet):
    return [('画外音 ', True), ('邀请 ', False), ('您', True),
            (' 担任“', False), (meet, True), ('”的', False), ('致辞嘉宾', True)]

def tail_summary_row(inviter, meet):
    return [(inviter + ' 教授', True), (' 邀请 ', False), ('您', True),
            (' 担任“', False), (meet, True), ('” 的', False), ('总结嘉宾', True)]

def build_task(t, row, meet):
    row = dict(row)
    if 'titles' in row:
        row['titles'] = [resolve(x) for x in row['titles']]
    if 'mandates' in row:
        row['mandates'] = [resolve(x) for x in row['mandates']]
    ty = row['type']
    if ty == 'drawn_invite':
        segs = drawn_invite_row(meet)
    elif ty == 'invite':
        segs = invite_row(row.get('inviter'), row['guest'],
                          row['titles'], row['role'])
    elif ty == 'disc_list':
        segs = disc_list_row(row['names'])
    elif ty == 'handoff':
        segs = handoff_row(row['next'], row['mandates'])
    elif ty == 'team_summary':
        segs = team_summary_row()
    elif ty == 'tail_summary':
        segs = tail_summary_row(row['inviter'], meet)
    else:
        raise RuntimeError('未知 row type: ' + ty)
    return {'time': t, 'paras': [segs]}

# ---------- 页管理 ----------
def clone_slide(prs, src, materialize=False):
    """克隆一页；materialize=True 时跨文件克隆（图片字节物化进目标包）。
    注意新页必须挂目标包自己的 layout，否则包非法、LibreOffice 拒载。"""
    blank = prs.slides.add_slide(prs.slides[0].slide_layout)
    cs_new = blank._element.find(qn('p:cSld'))
    cs_src = src._element.find(qn('p:cSld'))
    bg = cs_src.find(qn('p:bg'))
    if bg is not None:
        cs_new.insert(0, copy.deepcopy(bg))
    tree_new = blank.shapes._spTree
    for el in cs_src.find(qn('p:spTree')):
        tag = el.tag.split('}')[-1]
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            continue
        tree_new.append(copy.deepcopy(el))
    import tempfile
    R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    for el in blank._element.iter():
        for attr in (R + 'embed', R + 'id', R + 'link'):
            rid = el.get(attr)
            if not rid or rid not in src.part.rels:
                continue
            rel = src.part.rels[rid]
            tgt = rel.target_part
            try:
                if materialize and 'image' in getattr(tgt, 'content_type', ''):
                    ext = os.path.splitext(tgt.partname)[1].lstrip('.') or 'png'
                    tmpf = tempfile.NamedTemporaryFile(suffix='.' + ext, delete=False)
                    tmpf.write(tgt.blob)
                    tmpf.close()
                    _, new_rid = blank.part.get_or_add_image_part(tmpf.name)
                else:
                    new_rid = blank.part.relate_to(tgt, rel.reltype)
            except Exception:
                continue
            el.set(attr, new_rid)
    return blank

def delete_slide(prs, slide_obj):
    lst = prs.slides._sldIdLst
    elems = list(lst)
    cur = list(prs.slides)
    i = cur.index(slide_obj)
    prs.part.drop_rel(elems[i].get(qn('r:id')))
    lst.remove(elems[i])

def apply_order(prs, ordered):
    lst = prs.slides._sldIdLst
    elems = list(lst)
    cur = list(prs.slides)
    m = {id(cur[i]): elems[i] for i in range(len(cur))}
    for e in elems:
        lst.remove(e)
    for so in ordered:
        lst.append(m[id(so)])

# ---------- 各 kind 构建 ----------
def _finish(prs, letter, head_role):
    fix_salutation(prs, letter['name'])
    rewrite_intro(prs, head_role)
    global_replace(prs)
    outpath = os.path.join(CFG['output'], letter['file'])
    os.makedirs(os.path.dirname(outpath), exist_ok=True)
    prs.save(outpath)

def spec_chair(letter):
    prs = Presentation(BASE['chair'])
    meet = CFG['meeting']['intro']
    tasks = [build_task(t['time'], t['row'], meet) for t in letter['tasks']]
    fill_table(first_tables(prs)[0], tasks[:1])
    fill_table(first_tables(prs)[1], tasks[1:], header=False, renumber=False)
    _finish(prs, letter, None)

def spec_chairL(letter):
    out = Presentation(BASE['chair'])
    srchost = Presentation(BASE['host'])
    page = clone_slide(out, srchost.slides[1], materialize=True)
    rewrite_profile(page, letter['profile_expert'])
    gf = [sh for sh in out.slides[0].shapes if getattr(sh, 'has_table', False)]
    ROW_H = 648000
    meet = CFG['meeting']['intro']
    tasks = [build_task(t['time'], t['row'], meet) for t in letter['tasks']]
    n_main = len(letter['tasks']) - 1          # 主表行数（尾行属尾表）
    fill_table(gf[0].table._tbl, tasks[:n_main])
    gf[1].top = gf[1].top + (n_main - 1) * ROW_H   # 主表长高→尾表下移
    fill_table(gf[1].table._tbl, tasks[n_main:], header=False, renumber=False)
    tc0 = gf[1].table._tbl.findall(A('tr'))[0].findall(A('tc'))[0]
    t_el = list(tc0.iter(A('r')))[0].find(A('t'))
    if t_el is not None:
        t_el.text = str(n_main + 1)            # 尾表序号顺延
    _finish(out, letter, None)

def spec_host(letter):
    prs = Presentation(BASE['host'])
    sls = list(prs.slides)
    for i in sorted(letter.get('del_idx', []), reverse=True):
        delete_slide(prs, sls[i])
    live = list(prs.slides)
    for k, expert in enumerate(letter['profiles']):
        rewrite_profile(live[k + 1], expert)
    if letter.get('questions'):
        rewrite_topics(live[len(letter['profiles']) + 1], letter['questions'])
    if letter.get('next'):
        rewrite_profile(live[-1], letter['next'])
    meet = CFG['meeting']['intro']
    tasks = [build_task(t['time'], t['row'], meet) for t in letter['tasks']]
    fill_table(first_tables(prs)[0], tasks)
    _finish(prs, letter, '主持')

def spec_host6(letter):
    """6 张中间简介页：克隆 1 页补位"""
    prs = Presentation(BASE['host'])
    sls = list(prs.slides)
    clone = clone_slide(prs, sls[5])
    order = [sls[0], sls[1], sls[2], sls[3], sls[4], sls[5], clone, sls[6], sls[7]]
    for slide, expert in zip(order[1:7], letter['profiles']):
        rewrite_profile(slide, expert)
    rewrite_topics(order[7], letter['questions'])
    rewrite_profile(order[8], letter['next'])
    meet = CFG['meeting']['intro']
    tasks = [build_task(t['time'], t['row'], meet) for t in letter['tasks']]
    fill_table(first_tables(prs)[0], tasks)
    _finish(prs, letter, '主持')

def spec_hostL(letter):
    prs = Presentation(BASE['hostL'])
    sls = list(prs.slides)
    need_mid = len(letter['profiles'])
    mids, topics, trail = sls[1:5], sls[5], sls[6:10]
    order = [sls[0]] + mids[:]
    while len(order) - 1 < need_mid:           # 中间页不足4张时克隆补位
        c = clone_slide(prs, mids[-1])
        order.append(c)
    topics_i = len(order)
    order.append(topics)
    order += trail[:len(letter['trailing'])]
    for extra in trail[len(letter['trailing']):]:
        delete_slide(prs, extra)
    meet = CFG['meeting']['intro']
    for slide, expert in zip(order[1:1+need_mid], letter['profiles']):
        rewrite_profile(slide, expert)
    rewrite_topics(order[topics_i], letter['questions'])
    for slide, expert in zip(order[topics_i+1:], letter['trailing']):
        rewrite_profile(slide, expert)
    tasks = [build_task(t['time'], t['row'], meet) for t in letter['tasks']]
    fill_table(first_tables(prs)[0], tasks)
    rewrite_intro(prs, '主持')
    fix_salutation(prs, letter['name'])
    global_replace(prs)
    apply_order(prs, order)
    prs.save(os.path.join(CFG['output'], letter['file']))

def spec_simple(kind, letter):
    prs = Presentation(BASE[kind])
    if kind == 'disc':
        rewrite_topics(prs.slides[1], letter['questions'])
        rewrite_intro(prs, '讨论')
    else:
        rewrite_intro(prs, '演讲')
    meet = CFG['meeting']['intro']
    tasks = [build_task(t['time'], t['row'], meet) for t in letter['tasks']]
    fill_table(first_tables(prs)[0], tasks)
    fix_salutation(prs, letter['name'])
    global_replace(prs)
    prs.save(os.path.join(CFG['output'], letter['file']))

# ---------- 日程海报页（可选） ----------
def add_agenda_pages():
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as _PILImage
    img_path = CFG['agenda_img']
    iw, ih = _PILImage.open(img_path).size
    ar = iw / ih
    for fn in sorted(os.listdir(CFG['output'])):
        if not fn.endswith('.pptx') or fn.startswith('~$'):
            continue
        f = os.path.join(CFG['output'], fn)
        prs = Presentation(f)
        last = prs.slides[-1]
        covered = any(getattr(sh, 'shape_type', None) == MSO_SHAPE_TYPE.PICTURE
                      and (sh.width or 0) >= prs.slide_width * 0.85
                      and (sh.height or 0) >= prs.slide_height * 0.5
                      for sh in last.shapes)
        if covered:
            continue
        sl = prs.slides.add_slide(prs.slides[0].slide_layout)
        W, H = prs.slide_width, prs.slide_height
        bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bg.line.fill.background(); bg.shadow.inherit = False
        tw, th = W, int(W / ar)
        if th > H:
            th = H; tw = int(H * ar)
        sl.shapes.add_picture(img_path, (W - tw) // 2, (H - th) // 2, width=tw, height=th)
        prs.save(f)

# ---------- 配置装载 ----------
def resolve(x):
    """把 "@pos/@neg/@her2/@case/@kras..." 引用解析为实际文案"""
    if isinstance(x, str) and x.startswith('@'):
        key = x[1:]
        qs = CFG['questions']
        if key in qs:
            return qs[key]
        return CFG['titles'][key]
    return x

def load_config(path):
    global CFG, BASE, BIOS
    CFG = json.load(open(path, encoding='utf-8'))
    BASE = {k: os.path.join(CFG['template_dir'], v) for k, v in CFG['base'].items()}
    BIOS = {b['name']: b for b in json.load(open(CFG['bios'], encoding='utf-8'))}
    if CFG.get('bios_hospital_override'):
        for n, h in CFG['bios_hospital_override'].items():
            if n in BIOS:
                BIOS[n]['hospital'] = h
    os.makedirs(CFG['output'], exist_ok=True)
    # 捐体页脚里的旧日期/地点标记（用于全局替换识别）
    CFG['_old_hotel'] = CFG.get('old_hotel_mark', '佛山东平保利洲际酒店')
    CFG['_old_date_mark'] = CFG.get('old_date_mark', '会议日期：2026年8月2日')

def build_all():
    meet = CFG['meeting']['intro']
    for L in CFG['letters']:
        L = dict(L)
        L['file'] = '%s-%s-沟通函.pptx' % (CFG['prefix'][L['name']], L['name'])
        if L.get('questions'):
            L['questions'] = resolve(L['questions'])
        k = L['kind']
        if k in ('spk', 'disc') and 'tasks' not in L:
            role = '演讲嘉宾' if k == 'spk' else '学术讨论嘉宾'
            L['tasks'] = [{ 'time': L['time'],
                            'row': { 'type': 'invite', 'inviter': L['chair'],
                                     'guest': '您', 'titles': L['titles'], 'role': role } }]
        fn = {'chair': spec_chair, 'chairL': spec_chairL, 'host': spec_host,
              'host6': spec_host6, 'hostL': spec_hostL,
              'spk': lambda x: spec_simple('spk', x),
              'disc': lambda x: spec_simple('disc', x)}[k]
        done = False
        for _attempt in range(3):
            try:
                fn(L)
                done = True
                break
            except MemoryError:
                import gc; gc.collect()
        print(('OK  ' if done else 'ERR MemoryError ') + L['file'])
        if not done:
            sys.exit(1)
    if CFG.get('agenda_img'):
        add_agenda_pages()
        print('agenda pages appended')
    print('%d/%d generated -> %s' % (len(CFG['letters']), len(CFG['letters']), CFG['output']))

def fmt_time(t):
    def part(x):
        h, m = x.strip().split(':')
        return '%02d:%s' % (int(h), m)
    if '-' in t:
        a, b = t.split('-', 1)
        return '%s-%s' % (part(a), part(b))
    return part(t)

if __name__ == '__main__':
    load_config(sys.argv[1])
    build_all()
