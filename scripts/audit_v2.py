# -*- coding: utf-8 -*-
"""V2 沟通函审计：校验点由站点配置自动派生。
用法: python audit_v2.py <station.json>"""
import io, os, sys, json, hashlib, glob, re
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import build_v2 as bv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

AGENDA = 1 if bv.CFG.get('agenda_img') else 0
fails = []

def chk(cond, msg):
    if not cond:
        fails.append(msg)

def slide_text(sl):
    parts = []
    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(sh.shapes); continue
            if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
                parts.append(sh.text_frame.text)
            if getattr(sh, 'has_table', False) and sh.has_table:
                for r in sh.table.rows:
                    parts.append(' | '.join(c.text for c in r.cells))
    walk(sl.shapes)
    return '\n'.join(parts)

def row_text(row):
    if 'names' in row:
        return '、'.join(row['names']) + ' 教授'
    return ''.join(seg for para in [bv.build_task('', row, bv.CFG['meeting']['intro'])['paras'][0]] for seg, _ in para)

def base_slides(L):
    k = L['kind']
    if k == 'chair': return 1
    if k == 'chairL': return 2
    if k == 'spk': return 1
    if k == 'disc': return 2
    if k == 'hostL': return 1 + len(L['profiles']) + 1 + len(L['trailing'])
    return 1 + len(L['profiles']) + (1 if L.get('questions') else 0) + (1 if L.get('next') else 0)

bv.load_config(sys.argv[1])
meet = bv.CFG['meeting']
BASE_HASH = set()
for f in glob.glob(os.path.join(bv.CFG['template_dir'], '*.pptx')):
    if os.path.basename(f).startswith('~$'): continue
    prs = Presentation(f)
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.image is not None:
                BASE_HASH.add(hashlib.md5(sh.image.blob).hexdigest())

date_key = '%s（%s）上午%s' % (meet['date'], meet['weekday'], meet['span'])
STALE = bv.CFG.get('stale_names', [])

for L in bv.CFG['letters']:
    nm = L['name']
    path = os.path.join(bv.CFG['output'],
                        '%s-%s-沟通函.pptx' % (bv.CFG['prefix'][nm], nm))
    if not os.path.exists(path):
        fails.append(f'{nm}: 文件不存在'); continue
    prs = Presentation(path)
    texts = [slide_text(s) for s in prs.slides]
    joined = '\n'.join(texts)
    exp = base_slides(L) + AGENDA
    chk(len(texts) == exp, f'{nm}: 页数 {len(texts)}!={exp}')
    chk(('尊敬的' in texts[0] and nm in texts[0].split('尊敬的')[1][:40]), f'{nm}: 称呼缺名字')
    chk(date_key.replace(' ', '') in texts[0].replace(' ', ''), f'{nm}: 开场白日期未更新')
    chk(meet['hotel'] in joined, f'{nm}: 会议地点未更新')
    for bad in bv.CFG.get('stale_marks', []):
        chk(bad not in joined, f'{nm}: 残留 {bad}')
    for bad in STALE:
        chk(bad not in joined, f'{nm}: 残留 {bad}')
    if re.search(r'杨雪(?!琴)', joined):
        fails.append(f'{nm}: 残留 杨雪')
    for t in L.get('tasks', []):
        chk(t['time'] in texts[0], f'{nm}: 缺时间行 {t["time"]}')
        rt = row_text(t['row']).replace(' ', '')
        chk(rt in texts[0].replace(' ', ''), f'{nm}: 缺任务文案“{rt[:24]}”')
    if L.get('questions'):
        qs = bv.resolve(L['questions'])
        qk = qs[0][:12]
        chk(qk in joined, f'{nm}: 话题页缺“{qk}…”')
    people = L.get('profiles', []) + L.get('trailing', []) \
        + ([L['next']] if L.get('next') else []) \
        + ([L['profile_expert']] if L.get('profile_expert') else [])
    for pe in people:
        chk(any(pe in t2.replace(' ', '') and '教授' in t2 for t2 in texts), f'{nm}: 简介页缺 {pe}')
    for i, sl in enumerate(prs.slides):
        def walkp(shapes):
            for sh in shapes:
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walkp(sh.shapes); continue
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.image is not None:
                    h = hashlib.md5(sh.image.blob).hexdigest()
                    chk(h not in BASE_HASH or sh.image.size[0] > 2000, f'{nm}: s{i+1} 使用模板原图')
        walkp(sl.shapes)

# chairL 专项：尾表序号顺延 + 位移
cl = [L for L in bv.CFG['letters'] if L['kind'] == 'chairL']
for L in cl:
    path = os.path.join(bv.CFG['output'],
                        '%s-%s-沟通函.pptx' % (bv.CFG['prefix'][L['name']], L['name']))
    prs = Presentation(path)
    gfs = [sh for sh in prs.slides[0].shapes if getattr(sh, 'has_table', False)]
    if len(gfs) != 2:
        fails.append(f"{L['name']}: 应有两张表"); continue
    want = str(len(L['tasks']))
    chk(gfs[1].table.rows[0].cells[0].text.strip() == want, f"{L['name']}: 尾表序号应为{want}")
    from pptx.util import Emu
    gap = Emu(gfs[1].top - gfs[0].top).inches
    chk(gap > 1.7, f"{L['name']}: 尾表位移不足 {gap:.2f}in")

print('FAILS:', len(fails))
for f_ in fails:
    print(' -', f_)
sys.exit(1 if fails else 0)
