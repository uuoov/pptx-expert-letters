# -*- coding: utf-8 -*-
"""简历提取：串场PPT(按姓名匹配批量) / 单份docx·pptx简历 -> bios.json

用法:
  python extract_bios.py --deck 串场.pptx --names 周彩存,宋勇,黄诚 --out bios.json [--append]
  python extract_bios.py --file 陈群简介.docx --name 陈群 --out bios.json [--append]
  python extract_bios.py --file 王长利.pptx --name 王长利 --out bios.json [--append]
默认覆盖 out 文件；--append 合并进已有文件。
"""
import sys, io, os, json, re, zipfile, argparse

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

def para_text_with_br(p_el):
    """按顺序拼接 run 文本，<a:br>/<w:br> 记为换行——避免软换行丢字"""
    parts = []
    for node in p_el.iter():
        if node.tag.endswith('}t') and node.text:
            parts.append(node.text)
        elif node.tag.endswith('}br'):
            parts.append('\n')
    return ''.join(parts)

def from_deck(deck, names):
    from pptx import Presentation
    prs = Presentation(deck)
    bios = {}
    for i, s in enumerate(prs.slides, 1):
        texts = []
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append((t, [para_text_with_br(p._p).strip()
                                      for p in sh.text_frame.paragraphs
                                      if para_text_with_br(p._p).strip()]))
        who = None
        for t, _ in texts:
            tt = re.sub(r'[\s\u3000]', '', t)
            if tt in names:
                who = tt
                break
        if not who or who in bios:
            continue
        best = max(texts, key=lambda x: len(x[1]))
        if len(best[1]) >= 5:
            bios[who] = best[1]
    return bios

def from_resume(path):
    if path.lower().endswith('.docx'):
        with zipfile.ZipFile(path) as z:
            xml = z.read('word/document.xml').decode('utf-8')
        lines = []
        for p in re.findall(r'<w:p[ >].*?</w:p>', xml, re.S):
            t = ''.join(re.findall(r'<w:t[^>]*>([^<]*)</w:t>', p)).strip()
            if t:
                lines.append(t)
        return lines
    from pptx import Presentation
    from pptx.oxml.ns import qn
    prs = Presentation(path)
    lines = []
    for s in prs.slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame._txBody.findall(qn('a:p')):
                t = para_text_with_br(p).strip()
                if t:
                    lines.extend(x.strip() for x in t.split('\n') if x.strip())
    return lines

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--deck')
    ap.add_argument('--file')
    ap.add_argument('--names')
    ap.add_argument('--name')
    ap.add_argument('--out', required=True)
    ap.add_argument('--append', action='store_true')
    a = ap.parse_args()
    out = {}
    if a.append and os.path.exists(a.out):
        out = json.load(open(a.out, encoding='utf-8'))
    if a.deck:
        names = [n.strip() for n in a.names.split(',')]
        got = from_deck(a.deck, names)
        out.update(got)
        miss = [n for n in names if n not in got]
        print('串场抽取 %d 人；未找到: %s' % (len(got), miss or '无'))
    if a.file:
        lines = from_resume(a.file)
        out[a.name] = lines
        print('%s: %d 行' % (a.name, len(lines)))
    json.dump(out, open(a.out, 'w', encoding='utf-8'), ensure_ascii=False, indent=1)
    print('已写入 %s (共 %d 人)' % (a.out, len(out)))

if __name__ == '__main__':
    main()
