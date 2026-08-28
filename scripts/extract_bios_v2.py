# -*- coding: utf-8 -*-
"""V2 简历提取：目录驱动，输出 bios JSON + 裁剪好人像。

用法:
    python extract_bios_v2.py --res <简历目录> --out bios.json --port portraits/ \
        [--conv converted/] [--soffice "C:\\Program Files\\LibreOffice\\program\\soffice.exe"] \
        [--render-scans scans/] [--supplement supplement.json]

规则与坑（详见 references/v2-station-playbook.md）:
- 文件名按 `NN-角色-姓名-医院-简历.ext` 解析；同名取首个（如 01/11 两组主席）。
- .ppt/.doc 先 soffice 转换到 --conv（转换产物与原件分目录，按精确扩展名取用）。
- pptx 文本取"画布内文字最多的一页"（过滤离画布形状——那是别人的内容），
  递归进 GROUP，按 <a:br> 拆行；docx 直接解析 word/document.xml（python-docx
  对部分 WPS 文件抛 KeyError）；pdf 用 PyMuPDF。
- 文本为 0 = 扫描件：首页渲染 PNG 到 --render-scans-dir 供人工转写，
  转写结果写 supplement.json（{"姓名": {"lines":[...], "photo_frac":[l,t,r,b], "photo_png":"x.png"}}）
  后用 --supplement 合并。
- 照片取画布内最大显示面积图片，居中裁剪到 0.739（1.98×2.68in 版式框比例）。
- 文本清理：邮箱/自我署名/大会主席头/超长句（优先逗号处断行）；
  获奖清单整段裁剪、硬上限 17 行，截断名单打印。
"""
import io, os, re, json, sys, zipfile, argparse, subprocess
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
FRAME_AR = 1.98 / 2.68
TRIM_MARKERS = ['获奖情况', '主要获奖', '奖励情况']

def para_lines(tf):
    out = []
    for pe in tf._txBody.iter(A + 'p'):
        buf = []
        for node in pe.iter():
            tag = node.tag.split('}')[-1]
            if tag == 't': buf.append(node.text or '')
            elif tag == 'br': buf.append('\n')
        for piece in ''.join(buf).split('\n'):
            piece = re.sub(r'\s{2,}', ' ', piece).strip()
            if piece: out.append(piece)
    return out

def canvas_shapes(sl, SW):
    def walk(shapes):
        for sh in shapes:
            L = Emu(sh.left or 0).inches
            W = Emu(sh.width or 0).inches
            if L < -0.05 or L + W > SW + 0.05:
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield sh
                yield from walk(sh.shapes)
            else:
                yield sh
    yield from walk(sl.shapes)

def extract_pptx(path):
    prs = Presentation(path)
    SW = Emu(prs.slide_width).inches
    best = (-1, None)
    for sl in prs.slides:
        n = sum(len(sh.text_frame.text) for sh in canvas_shapes(sl, SW) if sh.has_text_frame)
        if n > best[0]: best = (n, sl)
    sl = best[1]
    tbs = [sh for sh in canvas_shapes(sl, SW) if sh.has_text_frame and sh.text_frame.text.strip()
           and not str(sh.shape_type).startswith('GROUP')]
    lines = para_lines(max(tbs, key=lambda s: len(s.text_frame.text)).text_frame) if tbs else []
    pic, area = None, 0
    for sl2 in prs.slides:
        for sh in canvas_shapes(sl2, SW):
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.image is not None:
                a = (sh.width or 0) * (sh.height or 0)
                if a > area:
                    try: sh.image.size
                    except Exception: continue
                    area, pic = a, sh.image.blob
    return lines, pic

def extract_docx(path):
    import xml.etree.ElementTree as ET
    W = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
    z = zipfile.ZipFile(path)
    root = ET.fromstring(z.read('word/document.xml'))
    lines = []
    for p in root.iter(W + 'p'):
        t = re.sub(r'\s{2,}', ' ', ''.join(n.text or '' for n in p.iter(W + 't'))).strip()
        if t: lines.append(t)
    blob, big = None, 0
    for i in z.infolist():
        if i.filename.startswith('word/media/'):
            try:
                im = Image.open(io.BytesIO(z.read(i.filename))); w, h = im.size
            except Exception: continue
            if w * h > big: big, blob = w * h, z.read(i.filename)
    return lines, blob

def extract_pdf(path):
    import fitz
    doc = fitz.open(path)
    blk = []
    for pg in doc:
        d = pg.get_text('dict')
        for b in d['blocks']:
            if b['type'] == 0:
                for ln in b['lines']:
                    t = ''.join(sp['text'] for sp in ln['spans']).strip()
                    if t: blk.append((b['bbox'][1], ln['bbox'][0], t))
    blk.sort()
    blob, big = None, 0
    for pi in range(len(doc)):
        for img in doc[pi].get_images(full=True):
            try:
                dd = doc.extract_image(img[0]); bl = dd['image']
                im = Image.open(io.BytesIO(bl)); w, h = im.size
            except Exception: continue
            if w * h > big: big, blob = w * h, bl
    return [t for _, _, t in blk], blob

def split_long(lines, maxlen=80):
    res = []
    for l in lines:
        if len(l) <= maxlen: res.append(l); continue
        buf = ''
        for p in re.split(r'[;;。。](?=\S)', l):
            p = p.strip()
            if not p: continue
            cand = (buf + '；' + p) if buf else p
            if len(cand) <= maxlen: buf = cand
            else:
                if buf: res.append(buf)
                while len(p) > maxlen:
                    cut = p.rfind('，', int(maxlen * 0.6), maxlen)
                    cut = cut if cut > 0 else maxlen
                    res.append(p[:cut + 1] if p[cut] == '，' else p[:cut])
                    p = ' ' + p[cut + 1:]
                buf = p.strip()
        if buf: res.append(buf)
    return res

def clean(lines, name, hosp):
    ns, hs = re.sub(r'\s+', '', name), re.sub(r'\s+', '', hosp)
    out = []
    for l in lines:
        s = re.sub(r'\s+', '', l)
        if '@' in l or s == ns: continue
        if s.startswith('大会主席'): continue
        if s in (ns + '教授', ns + '教授' + hs): continue
        if s == hs and out: continue
        out.append(l)
    return out

def trim(lines, name):
    for i, l in enumerate(lines):
        if any(k in l for k in TRIM_MARKERS) and i >= 6:
            return lines[:i], f'{name}: 截掉获奖清单{len(lines)-i}行'
    if len(lines) > 17:
        return lines[:17], f'{name}: 末尾截断{len(lines)-17}行'
    return lines, None

def crop_save(blob_or_path, out_path):
    im = Image.open(blob_or_path if isinstance(blob_or_path, str) else io.BytesIO(blob_or_path)).convert('RGB')
    w, h = im.size
    if abs(w / h - FRAME_AR) > 0.02:
        if w / h > FRAME_AR:
            nw = int(h * FRAME_AR); x0 = (w - nw) // 2; im = im.crop((x0, 0, x0 + nw, h))
        else:
            nh = int(w / FRAME_AR); y0 = (h - nh) // 2; im = im.crop((0, y0, w, y0 + nh))
    im.save(out_path, 'JPEG', quality=88)

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--res', required=True)
    ap.add_argument('--out', required=True)
    ap.add_argument('--port', required=True)
    ap.add_argument('--conv', default=None)
    ap.add_argument('--soffice', default=r'C:\Program Files\LibreOffice\program\soffice.exe')
    ap.add_argument('--render-scans', default=None)
    ap.add_argument('--supplement', default=None)
    a = ap.parse_args()
    os.makedirs(a.port, exist_ok=True)
    if a.conv: os.makedirs(a.conv, exist_ok=True)

    experts = {}
    for f in os.listdir(a.res):
        m = re.match(r'(\d+)-(.+?)-(.+?)-(.+?)-简历\.(\w+)$', f)
        if not m: continue
        pref, role, name, hosp, ext = m.groups()
        if name in experts: continue
        experts[name] = {'prefix': pref, 'role': role, 'hospital': hosp,
                         'src': os.path.join(a.res, f), 'ext': ext.lower()}

    bios, notes, scans = [], [], []
    for name, info in experts.items():
        src, ext = info['src'], info['ext']
        if ext in ('ppt', 'doc'):
            want = '.pptx' if ext == 'ppt' else '.docx'
            conv = [c for c in os.listdir(a.conv)
                    if c.startswith(info['prefix'] + '-') and c.endswith(want)]
            assert conv, f'{name}: 未找到转换文件 {want}（先用 soffice 转换）'
            src = os.path.join(a.conv, conv[0]); ext = want.lstrip('.')
        lines = blob = None
        for _attempt in range(3):          # 32 位 Python 解大图偶发 MemoryError
            try:
                if ext == 'pptx': lines, blob = extract_pptx(src)
                elif ext == 'docx': lines, blob = extract_docx(src)
                else: lines, blob = extract_pdf(src)
                break
            except MemoryError:
                import gc; gc.collect()
                lines = blob = None
        lines = split_long(clean(lines, name, info['hospital']))
        lines, note = trim(lines, name)
        if note: notes.append(note)
        photo = None
        if not lines and a.render_scans:            # 扫描件：渲染等人工转写
            os.makedirs(a.render_scans, exist_ok=True)
            import fitz
            pdf = src
            if ext == 'pptx':
                import time
                pdf = os.path.join(a.render_scans, os.path.splitext(os.path.basename(src))[0] + '.pdf')
                for _ in range(3):
                    subprocess.run([a.soffice, '--headless', '--convert-to', 'pdf',
                                    '--outdir', a.render_scans, src], capture_output=True)
                    if os.path.exists(pdf): break
                    time.sleep(2)
                assert os.path.exists(pdf), f'{name}: soffice 转 PDF 失败'
            d = fitz.open(pdf)
            for i in range(min(len(d), 3)):
                d[i].get_pixmap(dpi=120).save(os.path.join(a.render_scans, f'{name}_p{i+1}.png'))
            scans.append(name)
        elif blob:
            photo = os.path.join(a.port, name + '.jpg')
            crop_save(blob, photo)
        bios.append({'name': name, 'prefix': info['prefix'], 'role': info['role'],
                     'hospital': info['hospital'], 'lines': lines, 'photo': photo})

    if a.supplement and os.path.exists(a.supplement):
        supp = json.load(open(a.supplement, encoding='utf-8'))
        by = {b['name']: b for b in bios}
        for name, s in supp.items():
            b = by[name]
            b['lines'] = s['lines']
            png = os.path.join(a.render_scans or '', s['photo_png'])
            crop_save(png, os.path.join(a.port, name + '.jpg'))
            b['photo'] = os.path.join(a.port, name + '.jpg')
            if name in scans: scans.remove(name)

    json.dump(bios, open(a.out, 'w', encoding='utf-8'), ensure_ascii=False, indent=1)
    print('extracted', len(bios), 'experts ->', a.out)
    print('截断:', notes)
    print('扫描件待人工转写:', scans or '无')

if __name__ == '__main__':
    main()
