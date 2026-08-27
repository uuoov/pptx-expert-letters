# -*- coding: utf-8 -*-
"""模板解剖：dump 沟通函模板的全部版式参数，产出"模板档案"。

用法:
    python analyze_template.py <模板.pptx> [输出.txt]
不指定输出则打印到 stdout。生成前必跑，所有生成参数以本输出为准。
"""
import sys, io, os, hashlib, zipfile

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

IN = 914400.0

def bg_hash(zf, part_path, rels_path):
    """取 slide/layout/master 背景引用的图片 md5（无背景返回 None）"""
    try:
        xml = zf.read(part_path).decode('utf-8', 'ignore')
    except KeyError:
        return None
    if '<p:bg>' not in xml:
        return None
    try:
        rels = zf.read(rels_path).decode('utf-8', 'ignore')
    except KeyError:
        return 'bg(no-rels)'
    m = None
    bg = xml[xml.find('<p:bg>'):xml.find('</p:bg>') + 7]
    for tok in bg.split('r:embed=')[1:]:
        rid = tok.split('"')[1]
        t = 'Target="../media/' + rid.split('rId')[-1]  # 粗略，不精确
        m = rid
        break
    # 精确: 从 rels 找 rid -> target
    import re
    mm = re.search(r'Id="%s"[^>]*Target="([^"]+)"' % m, rels) if m else None
    if mm:
        target = mm.group(1).replace('../', 'ppt/')
        try:
            return hashlib.md5(zf.read(target)).hexdigest() + ' (' + target + ')'
        except KeyError:
            return 'bg(missing %s)' % target
    return 'bg(%s)' % m

def walk(shapes, depth=0):
    for s in shapes:
        ind = '  ' * depth
        try:
            pos = ' pos=(%.2f,%.2f) size=(%.2fx%.2f)in' % (
                Emu(s.left).inches, Emu(s.top).inches, Emu(s.width).inches, Emu(s.height).inches)
        except Exception:
            pos = ''
        print('%s[shape %s] %s %s%s' % (ind, s.shape_id, s.shape_type, s.name, pos))
        if s.shape_type == 6:
            walk(s.shapes, depth + 1)
            continue
        if getattr(s, 'has_table', False):
            tbl = s.table._tbl
            grid = [int(int(g.get('w')) / IN * 100) / 100 for g in tbl.findall(qn('a:gridCol'))]
            print('%s  TABLE %dx%d cols=%s' % (ind, len(tbl.findall(qn('a:tr'))), len(grid), grid))
            for ri, tr in enumerate(tbl.findall(qn('a:tr'))):
                h = int(tr.get('h'))
                cells = []
                seen = set()
                for tc in tr.findall(qn('a:tc')):
                    if tc in seen:
                        continue
                    seen.add(tc)
                    txt = ''.join(t.text or '' for t in tc.iter(qn('a:t')))
                    cells.append(txt.strip().replace('\n', '⏎')[:28])
                print('%s  row%d h=%d EMU (%.2fin): %s' % (ind, ri, h, h / IN, cells))
            continue
        if not s.has_text_frame:
            continue
        for pi, p in enumerate(s.text_frame.paragraphs):
            pPr = p._p.find(qn('a:pPr'))
            ln = '-'
            if pPr is not None:
                l = pPr.find(qn('a:lnSpc'))
                if l is not None and l.find(qn('a:spcPct')) is not None:
                    ln = l.find(qn('a:spcPct')).get('val')
            runs = []
            for r in p.runs:
                sz = r.font.size.pt if r.font.size else '?'
                b = 'B' if r.font.bold else '-'
                u = 'U' if r.font.underline else '-'
                runs.append('[%r sz=%s %s%s]' % (r.text[:20], sz, b, u))
            if p.text.strip() or runs:
                print('%s  p%d lnSpc=%s: %s' % (ind, pi, ln, ' '.join(runs) if runs else repr(p.text[:40])))

def main():
    path = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else None
    buf = io.StringIO()
    old, sys.stdout = sys.stdout, buf
    prs = Presentation(path)
    print('=== 模板解剖: %s' % path)
    print('页面: %.2f x %.2f in, %d 页' % (Emu(prs.slide_width).inches, Emu(prs.slide_height).inches, len(prs.slides)))
    with zipfile.ZipFile(path) as z:
        media = [(n, hashlib.md5(z.read(n)).hexdigest()[:10], len(z.read(n))) for n in z.namelist() if n.startswith('ppt/media/') and z.read(n)]
        print('媒体文件:', media)
    for i, slide in enumerate(prs.slides, 1):
        print('\n---- 第 %d 页 (layout=%r) ----' % (i, slide.slide_layout.name))
        with zipfile.ZipFile(path) as z:
            n = i
            h = bg_hash(z, 'ppt/slides/slide%d.xml' % n, 'ppt/slides/_rels/slide%d.xml.rels' % n)
            hm = bg_hash(z, 'ppt/slideLayouts/slideLayout1.xml', 'ppt/slideLayouts/_rels/slideLayout1.xml.rels')
            print('  本页背景: %s | layout1 背景: %s' % (h or '无(用母版)', hm or '无'))
        walk(slide.shapes, 1)
    sys.stdout = old
    report = buf.getvalue()
    if out:
        open(out, 'w', encoding='utf-8').write(report)
        print('已写入', out)
    else:
        print(report)

if __name__ == '__main__':
    main()
