# -*- coding: utf-8 -*-
"""专家沟通函批量生成引擎：配置 JSON -> 整套 A4 信函式 PPTX。

用法:
    python build_letters.py <config.json>

配置 schema（真实样例见 examples/fuzhou_config.json）:
{
  "template":  "模板.pptx",              // 任一份现成沟通函（讨论/讲者函最通用）
  "profile_template": "主持函样例.pptx",  // 可选：其第2页作为简介页版式来源
  "output":    "out/",                   // 输出根目录
  "strip_letter_bg": true,               // 去信函页页级背景(露出母版信头)
  "meeting": {"intro": "……会议",          // 开场白中加粗的会议名
              "time": "会议时间：…", "place": "会议地点：…"},
  "bios": "bios.json",                   // 姓名 -> [简历行]
  "letters": [{
     "name": "陈群",                      // 收件人姓名（=称呼+文件名）
     "file": "主持/主持-陈群-沟通函.pptx", // 相对输出目录
     "kind": "host",   // chair|host|speaker|discuss|custom: 决定默认行高/空段/冒号
     "colon": false, "trailing_empty": false,      // 可覆盖 kind 默认
     "note": null,                         // "会议地点"后备注行(如 参会方式)
     "pian": null,                         // 篇章行文本
     "topics": null,                       // 讨论话题行，每项一段
     "tasks": [{"time": "11:10-11:20",
                "paras": [ [ ["由 ","n"],["张同梅 教授","b"] ],   // 段落=[加粗片段表]
                           [],                                    // 空段
                           [["请 您 …","n"]] ]}],
     "heights": [1390015],                 // 任务行高(EMU)，缺省用 kind 预设
     "profiles": [{"label":"演讲嘉宾","name":"胡志皇",          // 简历从 bios 取
                   "topic":["圆桌引导发言：","…"], "topic_label":"演讲主题"}]
  }]
}
加粗片段格式: ["文本", "b"或"n"]。
"""
import sys, io, os, gc, json, copy, math

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.oxml.ns import qn

# 杭州版实测行高预设(EMU)，可被 config.row_heights 覆盖
PRESETS = dict(
    chair_header=291893, chair_t1=687731, chair_t2=645160,
    host_header=321945, host_short=464023, host_mid=932180,
    host_long=1266825, host_xl=1390015,
    header=321945, pian=340360, task=1099185)
KIND_DEFAULTS = dict(
    chair=dict(colon=True, trailing_empty=False),
    host=dict(colon=False, trailing_empty=False),
    speaker=dict(colon=False, trailing_empty=True),
    discuss=dict(colon=False, trailing_empty=True),
    custom=dict(colon=False, trailing_empty=True))
KIND_HEIGHTS = dict(
    chair=['chair_t1', 'chair_t2'], host=['host_xl'],
    speaker=['task'], discuss=['task'], custom=['task'])

_TPL = {}
_PROFILE = None
BIOS = {}

def _template(path):
    if path not in _TPL:
        _TPL[path] = Presentation(path)
    return _TPL[path]

def _profile_slide(path):
    global _PROFILE
    if _PROFILE is None:
        _PROFILE = Presentation(path)
    return _PROFILE.slides[1]

def strip_slide_bg(prs):
    slide = prs.slides[0]
    cSld = slide._element.find(qn('p:cSld'))
    bg = cSld.find(qn('p:bg'))
    if bg is not None:
        for blip in bg.iter(qn('a:blip')):
            rid = blip.get(qn('r:embed'))
            if rid:
                try:
                    slide.part.drop_rel(rid)
                except Exception:
                    pass
        cSld.remove(bg)

def rebuild_para(p_el, segments, bold_rPr, norm_rPr):
    """segments: [(text, bold)]；链式 addnext 防倒序；rPr 来自模板抓取"""
    for r in p_el.findall(qn('a:r')):
        p_el.remove(r)
    for e in p_el.findall(qn('a:endParaRPr')):
        p_el.remove(e)
    prev = p_el.find(qn('a:pPr'))
    for text, bold in segments:
        r = p_el.makeelement(qn('a:r'), {})
        r.append(copy.deepcopy(bold_rPr if bold else norm_rPr))
        t = p_el.makeelement(qn('a:t'), {})
        t.text = text
        r.append(t)
        if prev is not None:
            prev.addnext(r)
        else:
            p_el.append(r)
        prev = r

def _seg_pairs(segs):
    return [(t, b == 'b') for t, b in segs]

def _set_lnspc(p_el, pct):
    pPr = p_el.find(qn('a:pPr'))
    if pPr is None:
        pPr = p_el.makeelement(qn('a:pPr'), {})
        p_el.insert(0, pPr)
    for old in pPr.findall(qn('a:lnSpc')):
        pPr.remove(old)
    ln = pPr.makeelement(qn('a:lnSpc'), {})
    sp = pPr.makeelement(qn('a:spcPct'), {})
    sp.set('val', str(pct))
    ln.append(sp)
    pPr.insert(0, ln)

def add_profile_page(prs, prof):
    """克隆简介页版式并替换文字。label/name/bio/topic/topic_label"""
    src_slide = _profile_slide(prs.cfg['profile_template'])
    new = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src_slide.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    label_box = bio_box = name_box = topic_box = topic_label_box = None
    for sh in new.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t.startswith('【'):
            if '主题' in t:
                topic_label_box = sh
            else:
                label_box = sh
        elif t.endswith('教授'):
            name_box = sh
        elif len(t) > 40:
            bio_box = sh
        else:
            topic_box = sh
    assert label_box and bio_box and name_box, '简介页版式框识别失败，检查 profile_template'
    p = label_box.text_frame.paragraphs[0]
    p.runs[0].text = '【%s】：' % prof['label']
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)
    p = name_box.text_frame.paragraphs[0]
    p.runs[0].text = prof['name'] + '  教授'
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)
    lines = list(prof.get('bio') or BIOS.get(prof['name']) or ['【简历待补充】'])
    if len(lines) > 17:
        lines = lines[:17]
    pct = 150000 if len(lines) < 14 else (130000 if len(lines) < 16 else 120000)
    tf = bio_box.text_frame
    rPr = copy.deepcopy(tf.paragraphs[0].runs[0]._r.find(qn('a:rPr')))
    for para in list(tf.paragraphs):
        para._p.getparent().remove(para._p)
    for line in lines:
        p_el = tf._txBody.makeelement(qn('a:p'), {})
        _set_lnspc(p_el, pct)
        r = p_el.makeelement(qn('a:r'), {})
        r.append(copy.deepcopy(rPr))
        t = p_el.makeelement(qn('a:t'), {})
        t.text = line
        r.append(t)
        p_el.append(r)
        tf._txBody.append(p_el)
    topic = prof.get('topic')
    if topic is None:
        for box in (topic_label_box, topic_box):
            if box is not None:
                box._element.getparent().remove(box._element)
        return
    p = topic_label_box.text_frame.paragraphs[0]
    p.runs[0].text = '【%s】：' % prof.get('topic_label', '演讲主题')
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)
    tf = topic_box.text_frame
    rPr = copy.deepcopy(tf.paragraphs[0].runs[0]._r.find(qn('a:rPr')))
    for para in list(tf.paragraphs):
        para._p.getparent().remove(para._p)
    est = sum(max(1, math.ceil(len(l) / 20)) for l in topic)
    for line in topic:
        p_el = tf._txBody.makeelement(qn('a:p'), {})
        r = p_el.makeelement(qn('a:r'), {})
        rp = copy.deepcopy(rPr)
        if est > 3:
            rp.set('sz', '1800')     # 长话题防截断
        r.append(rp)
        t = p_el.makeelement(qn('a:t'), {})
        t.text = line
        r.append(t)
        p_el.append(r)
        tf._txBody.append(p_el)

def build_letter(cfg, letter):
    prs = Presentation(cfg['template'])
    prs.cfg = cfg
    if cfg.get('strip_letter_bg', True):
        strip_slide_bg(prs)
    slide = prs.slides[0]
    box = next(s for s in slide.shapes if s.has_text_frame and '尊敬的' in s.text_frame.text)
    tf = box.text_frame

    # 称呼：整段重建（模板姓名 run 常被拆开），rPr 抓自模板保持 20pt 加粗下划线
    name = letter['name']
    p0 = tf.paragraphs[0]
    norm14 = bold20 = None
    for r in p0.runs:
        rp = r._r.find(qn('a:rPr'))
        if r.font.size is not None and r.font.size.pt == 20 and bold20 is None:
            bold20 = rp
        if r.font.size is not None and r.font.size.pt == 14 and norm14 is None:
            norm14 = rp
    rebuild_para(p0._p, [('尊敬的 ', False), (name, True), (' 教授，', False)], bold20, norm14)

    def para_by(prefix):
        for p in tf.paragraphs:
            if p.text.strip().startswith(prefix):
                return p
        raise RuntimeError('未找到段落: ' + prefix)

    meet = cfg['meeting']
    p1 = para_by('您好')
    bold0 = norm0 = None
    for r in p1.runs:
        rp = r._r.find(qn('a:rPr'))
        if r.font.bold and bold0 is None:
            bold0 = rp
        if not r.font.bold and norm0 is None:
            norm0 = rp
    rebuild_para(p1._p, [
        ('       您好，非常荣幸邀请您参加', False),
        ('“' + meet['intro'] + '”会议', True),
        ('，我们将会议信息进行了整理，以便您更清晰的知晓会议内容及安排，希望为您会议期间的繁忙工作提供便捷：', False)],
        bold0, norm0)
    rb = para_by('会议时间').runs[0]._r.find(qn('a:rPr'))
    rebuild_para(para_by('会议时间')._p, [(meet['time'], True)], rb, rb)
    rebuild_para(para_by('会议地点')._p, [(meet['place'], True)], rb, rb)

    note = letter.get('note')
    if note:
        src_p = para_by('会议地点')._p
        new_p = copy.deepcopy(src_p)
        src_p.addnext(new_p)
        rebuild_para(new_p, [(note, True)], rb, rb)
        # 删"您好"后的空段补偿行数，防"任务环节"压表
        for p in tf.paragraphs:
            if p.text.strip().startswith('您好'):
                nxt = p._p.getnext()
                if nxt is not None and nxt.tag == qn('a:p'):
                    txt = ''.join(t.text or '' for t in nxt.iter(qn('a:t')))
                    if not txt.strip():
                        nxt.getparent().remove(nxt)
                break

    if letter.get('colon'):
        for r in para_by('您的任务环节').runs:
            r.text = '您的任务环节：'
            break
    if not letter.get('trailing_empty'):
        last = tf.paragraphs[-1]
        if not last.text.strip():
            last._p.getparent().remove(last._p)

    # ---- 任务表 ----
    tbl = next(s for s in slide.shapes if getattr(s, 'has_table', False)).table._tbl
    trs = tbl.findall(qn('a:tr'))
    header_tr, pian_tr, task_tr, topic_tr = (trs + [None, None, None, None])[:4]
    tc2 = task_tr.findall(qn('a:tc'))[2]
    bold_rPr = norm_rPr = None
    for rPr in tc2.iter(qn('a:rPr')):
        if rPr.get('b') == '1' and bold_rPr is None:
            bold_rPr = rPr
        if rPr.get('b') != '1' and norm_rPr is None:
            norm_rPr = rPr
    bold_rPr, norm_rPr = copy.deepcopy(bold_rPr), copy.deepcopy(norm_rPr)

    tasks = letter.get('tasks') or []
    task_trs = [task_tr]
    for _ in range(len(tasks) - 1):
        new = copy.deepcopy(task_tr)
        task_trs[-1].addnext(new)
        task_trs.append(new)
    final = [header_tr]
    if letter.get('pian'):
        final.append(pian_tr)
    final += task_trs
    if letter.get('topics'):
        final.append(topic_tr)
    for tr in tbl.findall(qn('a:tr')):
        tbl.remove(tr)
    for tr in final:
        tbl.append(tr)

    header_tr.set('h', str(letter.get('header_h') or PRESETS['header']))
    heights = letter.get('heights') or [PRESETS[k] for k in KIND_HEIGHTS[letter.get('kind', 'custom')]]
    for tr, hv in zip(task_trs, heights):
        tr.set('h', str(hv))
    if letter.get('pian') and pian_tr is not None:
        pian_tr.set('h', str(PRESETS['pian']))

    def fill_cell(tr, ci, groups):
        txBody = tr.findall(qn('a:tc'))[ci].find(qn('a:txBody'))
        paras = txBody.findall(qn('a:p'))
        while len(paras) < len(groups):
            txBody.append(copy.deepcopy(paras[-1]))
            paras = txBody.findall(qn('a:p'))
        for extra in paras[len(groups):]:
            txBody.remove(extra)
        paras = txBody.findall(qn('a:p'))
        for p_el, segs in zip(paras, groups):
            rebuild_para(p_el, _seg_pairs(segs) if segs else [], bold_rPr, norm_rPr)

    if letter.get('pian') and pian_tr is not None:
        fill_cell(pian_tr, 0, [[(letter['pian'], False)]])
    for i, (tr, task) in enumerate(zip(task_trs, tasks)):
        fill_cell(tr, 0, [[(str(i + 1), False)]])
        fill_cell(tr, 1, [[(task['time'], False)]])
        paras = task['paras']
        if paras and not isinstance(paras[0], list):
            paras = [paras]
        fill_cell(tr, 2, paras)
    if letter.get('topics') and topic_tr is not None:
        fill_cell(topic_tr, 2, [[(t, False)] for t in letter['topics']])

    for prof in letter.get('profiles') or []:
        add_profile_page(prs, prof)

    outpath = os.path.join(cfg['output'], *letter['file'].split('/'))
    os.makedirs(os.path.dirname(outpath), exist_ok=True)
    prs.save(outpath)

def main():
    cfg_path = sys.argv[1]
    cfg = json.load(open(cfg_path, encoding='utf-8'))
    if cfg.get('bios'):
        BIOS.update(json.load(open(cfg['bios'], encoding='utf-8')))
    if cfg.get('row_heights'):
        PRESETS.update(cfg['row_heights'])
    results = {}
    for letter in cfg['letters']:
        d = KIND_DEFAULTS[letter.get('kind', 'custom')]
        letter.setdefault('colon', d['colon'])
        letter.setdefault('trailing_empty', d['trailing_empty'])
        done = False
        for attempt in range(3):
            try:
                build_letter(cfg, letter)
                done = True
                break
            except MemoryError:
                gc.collect()
            except PermissionError as e:
                print('ERR %s: 文件被占用(请在 PowerPoint/WPS 中关闭后重跑) %s' % (letter['file'], e))
                break
        print(('OK  ' if done else 'ERR ') + letter['file'])
        results[letter['file']] = done
    ok = sum(results.values())
    print('')
    print('%d/%d generated' % (ok, len(results)))
    if ok < len(results):
        print('FAILED:', [k for k, v in results.items() if not v])
        sys.exit(1)

if __name__ == '__main__':
    main()
